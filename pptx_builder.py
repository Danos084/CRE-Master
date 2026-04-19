"""
Build a branded PowerPoint report from listings.json
Slides:
  1.   Cover - date, total counts (lease + sale)
  2.   Executive Summary - key stats for both types
  ---- FOR LEASE SECTION ----
  3.   Lease Summary (by type & top suburbs)
  4-N. One slide per suburb with lease listings
  ---- FOR SALE SECTION ----
  N+1. Sale Summary (by type & top suburbs)
  N+2. One slide per suburb with sale listings
"""

import json
from collections import defaultdict
from datetime import datetime, timedelta
from pathlib import Path

import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DATA_DIR = Path("data")
LISTINGS_FILE = DATA_DIR / "listings.json"
PPTX_OUT = DATA_DIR / "Bayside_Commercial_Report.pptx"

NAVY  = RGBColor(0x0D, 0x21, 0x37)
GOLD  = RGBColor(0xC9, 0xA8, 0x4C)
GREEN = RGBColor(0x1A, 0x6B, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGREY = RGBColor(0xF4, 0xF6, 0xFA)
DGREY = RGBColor(0x44, 0x44, 0x44)
MGREY = RGBColor(0x88, 0x99, 0xAA)


def _solid_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _add_textbox(slide, left, top, width, height, text,
                 font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
                 wrap=True, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def _cover_slide(prs, lease_total, sale_total, lease_new, sale_new):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    _solid_fill(bg, NAVY)
    bg.line.fill.background()
    bar = slide.shapes.add_shape(1, 0, Inches(3.6), prs.slide_width, Inches(0.08))
    _solid_fill(bar, GOLD)
    bar.line.fill.background()
    _add_textbox(slide, 0.6, 0.5, 9, 1.0, "PRECISION PROPERTY", 13, True, GOLD)
    _add_textbox(slide, 0.6, 1.2, 9, 1.4,
                 "Bayside & Logan Corridor\nCommercial Property Report",
                 36, True, WHITE, font_name="Georgia")
    _add_textbox(slide, 0.6, 3.8, 9, 0.5,
                 datetime.today().strftime("Week of %d %B %Y"), 14, False, GOLD)
    stats = [
        (str(lease_total), "LEASE LISTINGS"),
        (str(lease_new),   "NEW LEASES"),
        (str(sale_total),  "SALE LISTINGS"),
        (str(sale_new),    "NEW SALES"),
    ]
    for i, (val, label) in enumerate(stats):
        x = 0.5 + i * 2.3
        _add_textbox(slide, x, 4.5, 2.2, 0.8, val,  32, True, WHITE)
        _add_textbox(slide, x, 5.2, 2.2, 0.4, label, 9, False, GOLD)
    _add_textbox(slide, 0.6, 6.7, 9, 0.3,
                 "Source: commercialrealestate.com.au", 9, False, MGREY)


def _section_divider(prs, title, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    _solid_fill(bg, NAVY)
    bg.line.fill.background()
    bar = slide.shapes.add_shape(1, 0, Inches(3.2), prs.slide_width, Inches(0.1))
    _solid_fill(bar, accent)
    bar.line.fill.background()
    _add_textbox(slide, 0.6, 0.5, 9, 0.6, "PRECISION PROPERTY", 13, True, accent)
    _add_textbox(slide, 0.6, 1.8, 9, 1.5, title, 44, True, WHITE, font_name="Georgia")
    _add_textbox(slide, 0.6, 3.5, 9, 0.5,
                 datetime.today().strftime("%d %B %Y"), 14, False, accent)


def _summary_slide(prs, rows, title, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hdr = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    _solid_fill(hdr, NAVY)
    hdr.line.fill.background()
    _add_textbox(slide, 0.4, 0.2, 9, 0.7, title, 24, True, WHITE)
    type_counts = defaultdict(int)
    suburb_counts = defaultdict(int)
    for r in rows:
        type_counts[r.get("type") or "Unspecified"] += 1
        suburb_counts[r.get("suburb", "")] += 1
    top_types   = sorted(type_counts.items(),   key=lambda x: -x[1])[:6]
    top_suburbs = sorted(suburb_counts.items(), key=lambda x: -x[1])[:6]
    _add_textbox(slide, 0.4, 1.3, 4.5, 0.4, "By Property Type",      14, True, NAVY)
    for i, (t, c) in enumerate(top_types):
        _add_textbox(slide, 0.4, 1.8 + i*0.42, 4.5, 0.4, f"  {t}  -  {c}", 12, False, DGREY)
    _add_textbox(slide, 5.5, 1.3, 4.5, 0.4, "Top Suburbs by Volume", 14, True, NAVY)
    for i, (s, c) in enumerate(top_suburbs):
        _add_textbox(slide, 5.5, 1.8 + i*0.42, 4.5, 0.4, f"  {s}  -  {c} listings", 12, False, DGREY)
    line = slide.shapes.add_shape(1, 0, Inches(7.0), prs.slide_width, Inches(0.05))
    _solid_fill(line, accent)
    line.line.fill.background()
    _add_textbox(slide, 0.4, 7.1, 9, 0.3,
                 "Precision Property | precisionprop.com.au | 0432 203 354", 9, False, MGREY)


def _suburb_slide(prs, suburb, rows, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hdr = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    _solid_fill(hdr, NAVY)
    hdr.line.fill.background()
    _add_textbox(slide, 0.3, 0.1, 8.5, 0.7, suburb.upper(), 22, True, WHITE)
    _add_textbox(slide, 8.5, 0.2, 1.4, 0.5,
                 f"{len(rows)} listing{'s' if len(rows)>1 else ''}",
                 11, False, accent, PP_ALIGN.RIGHT)
    col_labels = ["Address", "Type", "Size", "Price / Rental", "Agent", "Link"]
    col_x      = [0.2,        2.9,    4.6,    5.6,              7.1,     8.6]
    col_w      = [2.6,        1.6,    0.9,    1.4,              1.4,     1.3]
    for label, x, w in zip(col_labels, col_x, col_w):
        cell = slide.shapes.add_shape(1, Inches(x), Inches(0.95), Inches(w), Inches(0.32))
        _solid_fill(cell, accent)
        cell.line.fill.background()
        _add_textbox(slide, x+0.02, 0.96, w-0.04, 0.30, label, 8, True, NAVY)
    row_h = 0.47
    max_rows = min(len(rows), 13)
    for ri, listing in enumerate(rows[:max_rows]):
        y = 1.30 + ri * row_h
        row_bg = slide.shapes.add_shape(1, 0, Inches(y), prs.slide_width, Inches(row_h-0.03))
        _solid_fill(row_bg, LGREY if ri % 2 == 0 else WHITE)
        row_bg.line.fill.background()
        price = listing.get("price_or_rental") or listing.get("asking_rental") or "POA"
        vals = [
            listing.get("address", "")[:50],
            listing.get("type", ""),
            listing.get("size", ""),
            price,
            listing.get("listing_agent", "")[:20],
            listing.get("link", "")[:35] or "-",
        ]
        for val, x, w in zip(vals, col_x, col_w):
            _add_textbox(slide, x+0.02, y+0.02, w-0.04, row_h-0.06,
                         str(val), 8, False, DGREY, wrap=True)
    if len(rows) > max_rows:
        _add_textbox(slide, 0.2, 1.30+max_rows*row_h, 9, 0.3,
                     f"  + {len(rows)-max_rows} more listings - see Excel for full data",
                     9, False, DGREY)
    line = slide.shapes.add_shape(1, 0, Inches(7.3), prs.slide_width, Inches(0.04))
    _solid_fill(line, accent)
    line.line.fill.background()


def build_pptx(listings: dict) -> Path:
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    all_rows = list(listings.values())
    cutoff   = (datetime.today() - timedelta(days=7)).strftime("%Y-%m-%d")
    lease_rows = [r for r in all_rows if r.get("listing_type", "Lease") == "Lease"]
    sale_rows  = [r for r in all_rows if r.get("listing_type") == "Sale"]
    lease_new  = sum(1 for r in lease_rows if r.get("first_seen", "") >= cutoff)
    sale_new   = sum(1 for r in sale_rows  if r.get("first_seen", "") >= cutoff)
    _cover_slide(prs, len(lease_rows), len(sale_rows), lease_new, sale_new)
    _section_divider(prs, "For Lease Listings", GOLD)
    _summary_slide(prs, lease_rows, "Lease - Executive Summary", GOLD)
    lease_by_suburb = defaultdict(list)
    for r in lease_rows:
        lease_by_suburb[r.get("suburb", "Unknown")].append(r)
    for suburb in sorted(lease_by_suburb):
        _suburb_slide(prs, suburb, lease_by_suburb[suburb], GOLD)
    _section_divider(prs, "For Sale Listings", GREEN)
    _summary_slide(prs, sale_rows, "Sale - Executive Summary", GREEN)
    sale_by_suburb = defaultdict(list)
    for r in sale_rows:
        sale_by_suburb[r.get("suburb", "Unknown")].append(r)
    for suburb in sorted(sale_by_suburb):
        _suburb_slide(prs, suburb, sale_by_suburb[suburb], GREEN)
    prs.save(PPTX_OUT)
    print(f"PowerPoint saved -> {PPTX_OUT}  ({len(prs.slides)} slides | Lease: {len(lease_rows)}, Sale: {len(sale_rows)})")
    return PPTX_OUT


if __name__ == "__main__":
    with open(LISTINGS_FILE) as f:
        data = json.load(f)
    build_pptx(data)
